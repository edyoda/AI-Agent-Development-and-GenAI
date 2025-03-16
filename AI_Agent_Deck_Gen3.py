from pptx import Presentation
from pptx.util import Inches
from openai import OpenAI
import json, os
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.linear_model import LinearRegression
from sklearn.feature_extraction.text import CountVectorizer
from sklearn.decomposition import LatentDirichletAllocation

# Initialize OpenAI client
client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

# Step 1: Generate Slide Content using GPT
def generate_slide_content(topic):
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that creates slide deck content."},
            {"role": "user", "content": f"Create a 3-slide presentation outline about {topic}. Each slide should have a title and 3 bullet points."},
        ],
        max_tokens=200,
    )
    return response.choices[0].message.content.strip()

# Step 2: Parse Generated Content into Slides
def parse_content_to_slides(content):
    slides_data = []
    sections = content.split("\n\n")  # Split into slides
    for section in sections:
        lines = section.split("\n")
        title = lines[0].replace("Slide Title: ", "").strip()
        bullets = [line.replace("- ", "").strip() for line in lines[1:]]
        slides_data.append({"title": title, "bullets": bullets})
    return slides_data

# Step 3: Create Slides using python-pptx
def create_presentation(slides_data, output_file="presentation.pptx"):
    prs = Presentation()
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
        title = slide.shapes.title
        title.text = slide_data["title"]
        
        content = slide.shapes.placeholders[1].text_frame
        for bullet in slide_data["bullets"]:
            p = content.add_paragraph()
            p.text = bullet
            p.level = 0  # Top-level bullet point
    
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

# Step 4: Collect Feedback (Rating and Comments)
def collect_feedback():
    feedback = input("Rate the presentation (1-5): ")
    comments = input("Any specific feedback (e.g., design, content)? ")
    return int(feedback), comments

# Step 5: Save Feedback (Rating and Comments)
def save_feedback(topic, content, feedback, comments):
    feedback_data = {
        "topic": topic,
        "content": content,
        "feedback": feedback,
        "comments": comments,
    }
    with open("feedback_data.json", "a") as f:
        f.write(json.dumps(feedback_data) + "\n")

# Step 6: Analyze Comments Using Topic Modeling (LDA)
def analyze_comments():
    comments = []
    with open("feedback_data.json", "r") as f:
        for line in f:
            data = json.loads(line)
            comments.append(data["comments"])
    
    # Skip analysis if there are fewer than 2 feedback documents
    if len(comments) < 2:
        print("Not enough feedback data to analyze. Need at least 2 feedback entries.")
        return None
    
    # Adjust min_df and max_df based on the number of documents
    min_df = 1  # Allow terms that appear in at least 1 document
    max_df = 0.95  # Ignore terms that appear in more than 95% of documents
    
    # Use CountVectorizer to convert comments into a bag of words
    vectorizer = CountVectorizer(max_df=max_df, min_df=min_df, stop_words="english")
    comment_counts = vectorizer.fit_transform(comments)

    #print(comment_counts)
    
    # Apply Latent Dirichlet Allocation (LDA) to extract topics
    lda = LatentDirichletAllocation(n_components=min(3, len(comments)), random_state=42)
    lda.fit(comment_counts)
    
    # Extract and return the top words for each topic
    feature_names = vectorizer.get_feature_names_out()
    topics = []
    for topic_idx, topic in enumerate(lda.components_):
        top_words = [feature_names[i] for i in topic.argsort()[:-6:-1]]
        topics.append(top_words)
    
    return topics



# Step 7: Incorporate Feedback into Content Generation
def improve_content_with_feedback(topic, comments_analysis):
    # Use the analyzed comments to refine the prompt
    feedback_prompt = " ".join([f"Avoid {word}." if "less" in word else f"Include {word}." for topic_words in comments_analysis for word in topic_words])
    
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that creates slide deck content."},
            {"role": "user", "content": f"Create a 3-slide presentation outline about {topic}. Each slide should have a title and 3 bullet points. {feedback_prompt}"},
        ],
        max_tokens=200,
    )
    return response.choices[0].message.content.strip()

# Step 8: Main Function
def main():
    # Load or analyze feedback data
    try:
        comments_analysis = analyze_comments()
        print("Analyzed feedback topics:", comments_analysis)
    except:
        print("No feedback data found. Starting without learning.")
        comments_analysis = None
    
    # User input
    topic = input("Enter the topic for your presentation: ")
    
    # Generate content (with feedback if available)
    if comments_analysis:
        content = improve_content_with_feedback(topic, comments_analysis)
    else:
        content = generate_slide_content(topic)
    
    print("Generated Content:\n", content)
    
    # Parse content into slides
    slides_data = parse_content_to_slides(content)
    
    # Create presentation
    create_presentation(slides_data)
    
    # Collect feedback
    feedback, comments = collect_feedback()
    save_feedback(topic, content, feedback, comments)
    print("Thank you for your feedback!")

# Run the program
if __name__ == "__main__":
    main()
