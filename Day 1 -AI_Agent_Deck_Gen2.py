from pptx import Presentation
from pptx.util import Inches
from openai import OpenAI
import json, os
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.linear_model import LinearRegression

# Initialize OpenAI client
client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

# Step 1: Generate Slide Content using GPT
def generate_slide_content(num_slides,topic):
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that creates slide deck content."},
            {"role": "user", "content": f"Create a {num_slides}-slide presentation outline about {topic}. Each slide should have a title and 3 bullet points."},
        ],
        max_tokens=2000,
    )
    return response.choices[0].message.content.strip()

# Step 2: Parse Generated Content into Slides
def parse_content_to_slides(content):
    slides_data = []
    sections = content.split("\n\n")  # Split into slides
    for section in sections:
        lines = section.split("\n")
        title = lines[0].replace("Slide Title: ", "").strip()
        bullets = [line.replace("- ", "").strip() for line in lines[1:]]
        slides_data.append({"title": title, "bullets": bullets})
    return slides_data

# Step 3: Create Slides using python-pptx
def create_presentation(slides_data, output_file="presentation.pptx"):
    prs = Presentation()
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
        title = slide.shapes.title
        title.text = slide_data["title"]
        
        content = slide.shapes.placeholders[1].text_frame
        for bullet in slide_data["bullets"]:
            p = content.add_paragraph()
            p.text = bullet
            p.level = 0  # Top-level bullet point
    
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

# Step 4: Collect Feedback
def collect_feedback():
    feedback = input("Rate the presentation (1-5): ")
    comments = input("Any specific feedback (e.g., design, content)? ")
    return int(feedback), comments

# Step 5: Save Feedback
def save_feedback(topic, content, feedback, comments):
    feedback_data = {
        "topic": topic,
        "content": content,
        "feedback": feedback,
        "comments": comments,
    }
    with open("feedback_data.json", "a") as f:
        f.write(json.dumps(feedback_data) + "\n")

# Step 6: Train Feedback Model
def train_feedback_model():
    # Load feedback data
    topics = []
    contents = []
    ratings = []
    with open("feedback_data.json", "r") as f:
        for line in f:
            data = json.loads(line)
            topics.append(data["topic"])
            contents.append(data["content"])
            ratings.append(data["feedback"])
    
    # Combine topic and content for feature extraction
    combined_texts = [f"{topic} {content}" for topic, content in zip(topics, contents)]
    
    # Convert text to features using TF-IDF
    vectorizer = TfidfVectorizer()
    X = vectorizer.fit_transform(combined_texts)
    
    # Train a linear regression model
    model = LinearRegression()
    model.fit(X, ratings)
    
    return model, vectorizer

# Step 7: Predict Rating
def predict_rating(model, vectorizer, topic, content):
    combined_text = f"{topic} {content}"
    X = vectorizer.transform([combined_text])
    return model.predict(X)[0]

# Step 8: Improve Content Generation
def improve_content_generation(topic, num_slides, model, vectorizer):
    # Generate multiple content options
    content_options = [generate_slide_content(topic, num_slides) for _ in range(3)]
    
    # Predict ratings for each option
    ratings = [predict_rating(model, vectorizer, topic, content) for content in content_options]
    
    # Select the best option
    best_index = np.argmax(ratings)
    return content_options[best_index]

# Step 9: Main Function
def main():
    # Load or train the feedback model
    try:
        model, vectorizer = train_feedback_model()
    except FileNotFoundError:
        print("No feedback data found. Starting without learning.")
        model, vectorizer = None, None
    
    # User input
    topic = input("Enter the topic for your presentation: ")
    num_slides = input("Enter the number of slides : ")
    
    # Generate content (with learning if model is available)
    if model:
        content = improve_content_generation(topic, num_slides, model, vectorizer)
    else:
        content = generate_slide_content(topic, num_slides)
    
    print("Generated Content:\n", content)
    
    # Parse content into slides
    slides_data = parse_content_to_slides(content)
    
    # Create presentation
    create_presentation(slides_data)
    
    # Collect feedback
    feedback, comments = collect_feedback()
    save_feedback(topic, content, feedback, comments)
    print("Thank you for your feedback!")

# Run the program
if __name__ == "__main__":
    main()
